"""
신한카드 검색광고 게재보고 PPT 자동 생성기
============================================

사용법:
    python generate_report.py \
        --pc "<PC URL>" \
        --mo "<MO URL>" \
        --card "신한카드 Deep Once" \
        --out "출력파일.pptx"

또는 인자 없이 실행하면 스크립트 하단의 DEFAULT 값을 사용.

동작:
1. PC URL (card-search.naver.com/item)을 1920x1080 뷰포트로 열어
   상단부터 하단까지 모든 <details>를 펼친 상태로 viewport-height 단위로
   잘라 9~10장의 풀폭 캡처를 만든다 (누락 0).
2. MO URL (m-card-search.naver.com)도 동일 방식으로 모바일 캡처.
3. 샘플 PPT(Deep Once)와 동일한 레이아웃의 PPT로 조립.
"""
from __future__ import annotations
import argparse
import math
import os
import sys
import tempfile
import time as _time_mod
from pathlib import Path
from typing import List

# 모든 print 를 즉시 flush (백그라운드 출력 확인용)
try:
    sys.stdout.reconfigure(line_buffering=True)
    sys.stderr.reconfigure(line_buffering=True)
except Exception:
    pass


def _log(msg: str) -> None:
    print(f"[{_time_mod.strftime('%H:%M:%S')}] {msg}", flush=True)

from playwright.sync_api import sync_playwright, Page
from PIL import Image
import numpy as np
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# ---- Capture parameters ---------------------------------------------------

PC_VIEWPORT = (1920, 1080)
PC_SCREENSHOT_H = 989   # 샘플 PPT의 PC 이미지 높이와 동일하게 맞춤
PC_OVERLAP = 60         # 슬라이드 간 약간의 겹침 (누락 방지)

PC_TITLE_VIEWPORT = (801, 1311)  # 슬라이드 1 (vertical) 캡처용 좁은 뷰포트

MO_VIEWPORT = (404, 874)         # 샘플 PPT MO 이미지와 동일
MO_TITLE_DETAIL_VIEWPORT = (391, 1258)  # 샘플의 image12 비율
MO_OVERLAP = 50

MO_USER_AGENT = (
    "Mozilla/5.0 (Linux; Android 13; SM-S918N) "
    "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Mobile Safari/537.36"
)

# 네이버 검색결과 카드 OneBox 컨테이너 셀렉터 (PC/MO 동일)
SEARCH_ONEBOX_SELECTOR = "div.sp_cardsearch"
SEARCH_PC_VIEWPORT = (1280, 900)
SEARCH_MO_VIEWPORT = (412, 915)
HEADER_TEXT_RIGHT_SEARCH_PC = "네이버 신용카드 검색결과_PC"
HEADER_TEXT_RIGHT_SEARCH_MO = "네이버 신용카드 검색결과_MO"

# 신한카드 자체 안내 페이지 (PC/MO 동일 URL → 디바이스별로 다른 레이아웃)
HEADER_TEXT_RIGHT_LANDING_PC = "신한카드 안내 페이지_PC"
HEADER_TEXT_RIGHT_LANDING_MO = "신한카드 안내 페이지_MO"


# ---- Web capture ----------------------------------------------------------

def _open_all_details(page: Page) -> None:
    """페이지 내 모든 혜택 details/접이식 요소를 펼친다.
    네이버 카드 페이지는 click 이벤트에서 본문을 lazy-load 하므로
    반드시 summary 를 click() 해야 한다 (`.open = true` 만으론 본문이 안 채워짐).
    """
    # 1) 혜택 details (details.details: 렌탈/관리비/문화)
    benefit_summaries = page.locator('details.details > summary')
    cnt = benefit_summaries.count()
    for i in range(cnt):
        try:
            s = benefit_summaries.nth(i)
            s.scroll_into_view_if_needed(timeout=3000)
            s.click(timeout=3000)
            page.wait_for_timeout(350)
        except Exception as e:
            print(f"[WARN] benefit summary[{i}] click failed: {e}")

    # 2) 기타 details (annualFee_detail, baseRecord_detail 등): open 만 처리
    page.evaluate("""() => {
        document.querySelectorAll('details:not(.details)').forEach(d => { d.open = true; });
        document.querySelectorAll('[aria-expanded="false"]').forEach(el => {
            el.setAttribute('aria-expanded', 'true');
        });
    }""")
    # 클릭 후 페이지 리플로우 대기
    page.wait_for_timeout(600)


def _hide_floating_chrome(page: Page) -> None:
    """우측 상단 공유/맨위로 버튼 등 떠다니는 요소가 캡처를 가리지 않도록.
    (샘플 PPT에는 이 버튼들도 그대로 찍혀 있으므로 우선 숨기지 않음)
    """
    pass


def _full_page_height(page: Page) -> int:
    return page.evaluate("() => document.documentElement.scrollHeight")


def _capture_tall(page: Page, out_path: Path, base_w: int, max_h: int = 16000) -> int:
    """
    누락 없는 풀페이지 캡처:
    1) fixed/sticky → absolute (Playwright stitching 시 매 viewport 마다 중복 캡처되는 문제 해결)
    2) lazy-load 트리거 (페이지 하단까지 한 번 스크롤)
    3) 최상단으로 복귀 후 `full_page=True` 로 캡처
       — viewport 를 키우면 모바일 reflow 로 인해 페이지가 무한히 늘어나는
       문제(빈 슬라이스 생성)가 있어 viewport-resize 방식은 폐기.
    """
    page.wait_for_timeout(300)
    page.evaluate("() => window.scrollTo(0, 0)")
    page.wait_for_timeout(200)

    # fixed/sticky → absolute
    page.evaluate("""() => {
        document.querySelectorAll('*').forEach(el => {
            const cs = getComputedStyle(el);
            if (cs.position === 'fixed' || cs.position === 'sticky') {
                el.style.setProperty('position', 'absolute', 'important');
            }
        });
    }""")
    page.wait_for_timeout(200)

    # lazy-load 트리거: 페이지 끝까지 스크롤 → 이미지 로드 대기 → 다시 최상단
    page.evaluate("() => window.scrollTo(0, document.documentElement.scrollHeight)")
    page.wait_for_timeout(800)
    try:
        page.evaluate("""async () => {
            await Promise.all(
              Array.from(document.images).filter(img => !img.complete)
                .map(img => new Promise(res => { img.onload = img.onerror = res; }))
            );
        }""")
    except Exception:
        pass
    page.evaluate("() => window.scrollTo(0, 0)")
    page.wait_for_timeout(400)

    page_h = page.evaluate(
        "() => Math.max(document.documentElement.scrollHeight, document.body.scrollHeight)"
    )
    page.screenshot(path=str(out_path), full_page=True, animations="disabled")
    return page_h


def _slice_full_page(
    full_path: str, slice_h: int, overlap: int, out_dir: Path, prefix: str
) -> List[Path]:
    """긴 전체 페이지 스크린샷을 슬라이드 단위로 자른다 (겹침 포함)."""
    img = Image.open(full_path)
    W, H = img.size
    step = slice_h - overlap
    n = max(1, math.ceil((H - overlap) / step))
    paths: List[Path] = []
    for i in range(n):
        top = i * step
        bottom = min(top + slice_h, H)
        # 마지막 조각이 너무 짧으면 위로 당겨서 slice_h 확보
        if bottom - top < slice_h * 0.5 and i > 0:
            top = max(0, H - slice_h)
            bottom = H
        crop = img.crop((0, top, W, bottom))
        # 정확히 slice_h 가 안 되면 캔버스에 붙여서 패딩
        if crop.size[1] < slice_h:
            canvas = Image.new("RGB", (W, slice_h), (255, 255, 255))
            canvas.paste(crop, (0, 0))
            crop = canvas
        p = out_dir / f"{prefix}_{i + 1:02d}.png"
        crop.save(p, "PNG", optimize=True)
        paths.append(p)
        if bottom >= H:
            break
    return paths


def capture_pc(url: str, work: Path) -> dict:
    """PC 캡처: 1920×1080 뷰포트로 풀페이지 캡처 → 989px 단위 슬라이스 N장."""
    work.mkdir(parents=True, exist_ok=True)
    pc_full = work / "pc_full.png"
    _log("capture_pc: start")
    t0 = _time_mod.time()

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        ctx = browser.new_context(
            viewport={"width": PC_VIEWPORT[0], "height": PC_VIEWPORT[1]},
            device_scale_factor=1,
        )
        page = ctx.new_page()
        page.goto(url, wait_until="networkidle", timeout=60000)
        page.wait_for_timeout(1500)
        _open_all_details(page)
        page.wait_for_timeout(800)
        # 모든 이미지 로드 대기
        page.evaluate("""async () => {
            await Promise.all(
              Array.from(document.images)
                .filter(img => !img.complete)
                .map(img => new Promise(res => { img.onload = img.onerror = res; }))
            );
        }""")
        page.wait_for_timeout(500)
        height = _capture_tall(page, pc_full, PC_VIEWPORT[0])
        ctx.close()
        browser.close()

    slides = _slice_full_page(str(pc_full), PC_SCREENSHOT_H, PC_OVERLAP, work, "pc_slide")
    _log(f"capture_pc: done ({int(_time_mod.time()-t0)}s, {len(slides)} slides)")
    return {"slides": slides, "full": pc_full, "height": height}


def capture_mo(url: str, work: Path) -> dict:
    """MO 캡처: 404×874 모바일 뷰포트로 풀페이지 캡처 → 874px 단위 슬라이스 N장.
    PC 와 동일한 로직 — 별도의 검색결과/타이틀 캡처는 만들지 않음.
    """
    work.mkdir(parents=True, exist_ok=True)
    mo_full = work / "mo_full.png"
    _log("capture_mo: start")
    t0 = _time_mod.time()

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        ctx = browser.new_context(
            viewport={"width": MO_VIEWPORT[0], "height": MO_VIEWPORT[1]},
            user_agent=MO_USER_AGENT,
            is_mobile=True, has_touch=True, device_scale_factor=1,
        )
        page = ctx.new_page()
        page.goto(url, wait_until="networkidle", timeout=60000)
        page.wait_for_timeout(1500)
        _open_all_details(page)
        page.wait_for_timeout(800)
        page.evaluate("""async () => {
            await Promise.all(
              Array.from(document.images)
                .filter(img => !img.complete)
                .map(img => new Promise(res => { img.onload = img.onerror = res; }))
            );
        }""")
        page.wait_for_timeout(500)
        height = _capture_tall(page, mo_full, MO_VIEWPORT[0])
        ctx.close()
        browser.close()

    slides = _slice_full_page(str(mo_full), MO_VIEWPORT[1], MO_OVERLAP, work, "mo_slide")
    _log(f"capture_mo: done ({int(_time_mod.time()-t0)}s, {len(slides)} slides)")
    return {"slides": slides, "full": mo_full, "height": height}


# ---- 검색결과 OneBox 캡처 -------------------------------------------------

def capture_search_onebox(url: str, work: Path, is_mobile: bool = False) -> dict:
    """네이버 검색결과 페이지의 신용카드 OneBox 영역(div.sp_cardsearch)을
    통째로 캡처해 텍스트가 잘리지 않는 위치에서 2장으로 분할한다.

    - 카드 정보 + 연체이자율 + 법적고지 + "최종 업데이트 / 상품정보 오류신고"
      까지 한 컨테이너 안에 모두 있어 한 번에 캡처 가능.
    - 캡처본을 2등분하되, 단색 가로 띠(요소 간 공백)를 찾아 거기서 자르므로
      텍스트가 가운데에서 끊기지 않는다.

    Args:
        is_mobile: True 면 m.search.naver.com 모바일 검색결과로 가정해
                   모바일 UA/뷰포트 사용. False 면 PC.
    """
    work.mkdir(parents=True, exist_ok=True)
    prefix = "search_mo" if is_mobile else "search_pc"
    full = work / f"{prefix}_full.png"
    _log(f"capture_search_onebox[{'MO' if is_mobile else 'PC'}]: start")
    t0 = _time_mod.time()

    if is_mobile:
        vp = SEARCH_MO_VIEWPORT
        ctx_kw = dict(
            viewport={"width": vp[0], "height": vp[1]},
            user_agent=MO_USER_AGENT,
            is_mobile=True, has_touch=True, device_scale_factor=1,
        )
    else:
        vp = SEARCH_PC_VIEWPORT
        ctx_kw = dict(
            viewport={"width": vp[0], "height": vp[1]},
            device_scale_factor=1,
        )

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        ctx = browser.new_context(**ctx_kw)
        page = ctx.new_page()
        # 검색결과는 광고/트래커가 계속 떠서 networkidle 도 안 끝남
        # → commit (HTTP 응답 받자마자 진행) 후 OneBox 가 나타날 때까지만 짧게 대기
        _log(f"  search[{'MO' if is_mobile else 'PC'}]: goto...")
        page.goto(url, wait_until="commit", timeout=20000)
        _log(f"  search[{'MO' if is_mobile else 'PC'}]: goto returned, waiting for OneBox...")

        target = page.locator(SEARCH_ONEBOX_SELECTOR).first
        try:
            target.wait_for(state="visible", timeout=15000)
        except Exception:
            browser.close()
            raise RuntimeError(
                f"검색결과 페이지에서 카드 OneBox({SEARCH_ONEBOX_SELECTOR})를 찾지 못했습니다. "
                f"URL 이 카드명 검색결과({'모바일' if is_mobile else 'PC'})인지 확인해 주세요."
            )
        _log(f"  search[{'MO' if is_mobile else 'PC'}]: OneBox visible")

        target.scroll_into_view_if_needed()
        page.wait_for_timeout(400)
        # 이미지 로드 대기는 max 3초 (영원히 안 끝나는 트래커 이미지 회피)
        try:
            page.evaluate("""async () => {
                const promises = Array.from(document.images)
                  .filter(img => !img.complete)
                  .map(img => new Promise(res => {
                    img.onload = img.onerror = res;
                    setTimeout(res, 2500); // 개별 이미지 max 2.5초
                  }));
                await Promise.race([
                  Promise.all(promises),
                  new Promise(res => setTimeout(res, 3000)) // 전체 max 3초
                ]);
            }""")
        except Exception:
            pass
        page.wait_for_timeout(200)

        target.screenshot(path=str(full))
        browser.close()

    parts = _split_image_safely(full, work, prefix=f"{prefix}_part", n_parts=2)
    _log(f"capture_search_onebox[{'MO' if is_mobile else 'PC'}]: done ({int(_time_mod.time()-t0)}s)")
    return {"slides": parts, "full": full, "is_mobile": is_mobile}


# ---- 신한카드 자체 안내 페이지 캡처 ---------------------------------------

def capture_landing(url: str, work: Path, is_mobile: bool = False) -> dict:
    """신한카드 자체 카드 안내 페이지(shinhancard.com 등)를 풀페이지 캡처.
    - PC/MO 동일 URL 을 받지만 디바이스별 viewport/UA 로 따로 캡처.
    - 본문 캡처와 동일하게 풀폭 슬라이스 (PC 989px / MO 874px).
    - fixed/sticky 변환은 `_capture_tall` 내부에서 처리.
    """
    work.mkdir(parents=True, exist_ok=True)
    prefix = "landing_mo" if is_mobile else "landing_pc"
    full = work / f"{prefix}_full.png"
    _log(f"capture_landing[{'MO' if is_mobile else 'PC'}]: start")
    t0 = _time_mod.time()

    if is_mobile:
        vp = MO_VIEWPORT
        ctx_kw = dict(
            viewport={"width": vp[0], "height": vp[1]},
            user_agent=MO_USER_AGENT,
            is_mobile=True, has_touch=True, device_scale_factor=1,
        )
        slice_h = MO_VIEWPORT[1]
        overlap = MO_OVERLAP
    else:
        vp = PC_VIEWPORT
        ctx_kw = dict(
            viewport={"width": vp[0], "height": vp[1]},
            device_scale_factor=1,
        )
        slice_h = PC_SCREENSHOT_H
        overlap = PC_OVERLAP

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        ctx = browser.new_context(**ctx_kw)
        page = ctx.new_page()
        _log(f"  landing[{'MO' if is_mobile else 'PC'}]: goto...")
        page.goto(url, wait_until="domcontentloaded", timeout=30000)
        _log(f"  landing[{'MO' if is_mobile else 'PC'}]: goto returned")
        page.wait_for_timeout(1500)
        height = _capture_tall(page, full, vp[0])
        _log(f"  landing[{'MO' if is_mobile else 'PC'}]: captured (H={height})")
        ctx.close()
        browser.close()

    slides = _slice_full_page(str(full), slice_h, overlap, work, f"{prefix}_slide")
    _log(f"capture_landing[{'MO' if is_mobile else 'PC'}]: done ({int(_time_mod.time()-t0)}s, {len(slides)} slides)")
    return {"slides": slides, "full": full, "height": height, "is_mobile": is_mobile}


def _split_image_safely(
    img_path: Path, out_dir: Path, prefix: str, n_parts: int = 2,
) -> List[Path]:
    """이미지를 n_parts 등분하되, 텍스트 잘림 방지를 위해
    배경 단색 가로 띠(요소 간 공백) 위치를 분할점으로 사용.

    알고리즘:
      1) grayscale 변환 후 각 row 의 표준편차(std) 계산.
         - std 가 매우 낮은 row = 그 라인 전체가 단색 = 텍스트/그래픽 없음
      2) 목표 분할점 (전체 높이 × i/n) 의 ±15% 구간 내에서
         std 가 가장 낮은(연속으로 낮은) row 들의 중심을 분할점으로.
      3) 적절한 단색 구간이 없으면 fallback 으로 목표 분할점 그대로 사용.
    """
    img = Image.open(img_path).convert("RGB")
    W, H = img.size
    gray = np.array(img.convert("L"), dtype=np.float32)
    row_std = gray.std(axis=1)  # 각 row 의 표준편차 (W 픽셀에 대한)

    parts: List[Path] = []
    last_top = 0
    for i in range(n_parts - 1):
        target_y = int(H * (i + 1) / n_parts)
        # 분할 후보 범위 (이미지 중앙 근처)
        window = int(H * 0.15)
        s = max(last_top + 80, target_y - window)
        e = min(H - 80, target_y + window)
        if e <= s:
            split_y = target_y
        else:
            stds = row_std[s:e]
            # 단색에 가까운 row 들 (절대 임계값 + 상위 분위수 보정)
            threshold = max(stds.min() + 1.5, np.percentile(stds, 15))
            quiet = np.where(stds <= threshold)[0]
            if len(quiet) == 0:
                split_y = target_y
            else:
                # 연속된 단색 구간들 중 target_y 에 가장 가까운 구간의 중심
                # 그룹핑
                groups = []
                cur = [quiet[0]]
                for v in quiet[1:]:
                    if v - cur[-1] <= 2:
                        cur.append(v)
                    else:
                        groups.append(cur)
                        cur = [v]
                groups.append(cur)
                # 각 그룹의 중심 y (이미지 좌표)
                centers = [s + (g[0] + g[-1]) // 2 for g in groups]
                widths = [len(g) for g in groups]
                # target 과의 거리 + 두께(thicker = 더 안전한 공백) 가중 score
                best_idx = int(
                    np.argmin([abs(c - target_y) - widths[k] * 0.5 for k, c in enumerate(centers)])
                )
                split_y = centers[best_idx]

        crop = img.crop((0, last_top, W, split_y))
        p = out_dir / f"{prefix}_{i + 1:02d}.png"
        crop.save(p, "PNG", optimize=True)
        parts.append(p)
        last_top = split_y

    # 마지막 조각
    crop = img.crop((0, last_top, W, H))
    p = out_dir / f"{prefix}_{n_parts:02d}.png"
    crop.save(p, "PNG", optimize=True)
    parts.append(p)
    return parts


# ---- PPT building ---------------------------------------------------------

# 슬라이드 크기 (16:9, 12192000 x 6858000 EMU = 13.33" x 7.5")
SLIDE_W_EMU = Emu(12192000)
SLIDE_H_EMU = Emu(6858000)

TITLE_BLUE = RGBColor(0x44, 0x72, 0xC4)   # 좌상단 보조 텍스트 컬러
ACCENT_BLUE = RGBColor(0x5B, 0x9B, 0xD5)  # 라인 색

HEADER_TEXT_LEFT = "신한카드 멤버십 영업팀"   # 상단 좌측 (작은 청색)
HEADER_TEXT_RIGHT_PC = "네이버 신용카드 검색광고_PC"
HEADER_TEXT_RIGHT_MO = "네이버 신용카드 검색광고_MO"


def _add_header(slide, prs, right_text: str):
    """샘플 PPT 스타일 상단 헤더 (좌: 회사명, 우: PC/MO 표시)."""
    # 좌측 회사명
    tb = slide.shapes.add_textbox(Emu(234953), Emu(120000), Emu(6000000), Emu(420000))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = HEADER_TEXT_LEFT
    r = p.runs[0]
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = TITLE_BLUE

    # 우측 PC/MO 라벨
    tb2 = slide.shapes.add_textbox(Emu(8200000), Emu(120000), Emu(3800000), Emu(420000))
    tf2 = tb2.text_frame
    tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = 0
    p2 = tf2.paragraphs[0]
    p2.alignment = 2  # right
    p2.text = right_text
    r2 = p2.runs[0]
    r2.font.size = Pt(14)
    r2.font.bold = False
    r2.font.color.rgb = TITLE_BLUE


def _add_image_centered(slide, img_path: Path, top_emu: int = 700000,
                        max_w_emu: int = None, max_h_emu: int = None):
    """슬라이드 가운데 상단 정렬로 이미지 배치 (비율 유지)."""
    img = Image.open(img_path)
    w, h = img.size
    max_w_emu = max_w_emu or int(SLIDE_W_EMU * 0.95)
    max_h_emu = max_h_emu or (int(SLIDE_H_EMU) - top_emu - 200000)
    # EMU per pixel at 96 DPI: 9525
    px_to_emu = 9525
    target_w = w * px_to_emu
    target_h = h * px_to_emu
    scale = min(max_w_emu / target_w, max_h_emu / target_h, 1.0)
    final_w = int(target_w * scale)
    final_h = int(target_h * scale)
    left = (int(SLIDE_W_EMU) - final_w) // 2
    slide.shapes.add_picture(str(img_path), left, top_emu, width=final_w, height=final_h)


def _add_image_grid(slide, paths: List[Path], top_emu: int = 700000):
    """모바일 썸네일들을 한 줄로 가운데 정렬 배치."""
    if not paths:
        return
    gap_emu = Emu(150000)
    # 슬라이드 폭 95% 사용 / 길이만큼 분할
    avail_w = int(SLIDE_W_EMU * 0.95)
    avail_h = int(SLIDE_H_EMU) - top_emu - 200000
    # 첫 이미지로 비율 계산
    sample = Image.open(paths[0])
    sw, sh = sample.size
    # 각 카드 폭 추정 후 높이 맞추기
    n = len(paths)
    per_w_emu = (avail_w - int(gap_emu) * (n - 1)) // n
    # px→emu
    px_to_emu = 9525
    target_w = sw * px_to_emu
    target_h = sh * px_to_emu
    scale_w = per_w_emu / target_w
    scale_h = avail_h / target_h
    scale = min(scale_w, scale_h, 1.0)
    final_w = int(target_w * scale)
    final_h = int(target_h * scale)
    total_w = final_w * n + int(gap_emu) * (n - 1)
    left0 = (int(SLIDE_W_EMU) - total_w) // 2
    for i, pth in enumerate(paths):
        left = left0 + i * (final_w + int(gap_emu))
        slide.shapes.add_picture(str(pth), left, top_emu, width=final_w, height=final_h)


def build_pptx(
    pc: dict, mo: dict, out_path: Path,
    card_name: str,
    search_pc: dict | None = None,
    search_mo: dict | None = None,
    landing_pc: dict | None = None,
    landing_mo: dict | None = None,
    slides_per_mo_grid: int = 5, **_ignored,
) -> Path:
    """PC, MO (+ 선택적 검색결과, 안내페이지) 캡처를 PPT 로 조립.

    슬라이드 순서:
      [PC 검색결과 OneBox 2장]   ← search_pc
      [PC 상세 페이지 슬라이스]
      [PC 안내 페이지 슬라이스]  ← landing_pc (PC 섹션 맨 끝)
      [MO 검색결과 OneBox 2장]   ← search_mo
      [MO 상세 페이지 5장 그리드]
      [MO 안내 페이지 5장 그리드] ← landing_mo (MO 섹션 맨 끝)
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU
    blank_layout = prs.slide_layouts[6]

    # ---- PC: (검색결과) → 상세 → 안내 ----
    if search_pc and search_pc.get("slides"):
        for slice_path in search_pc["slides"]:
            s = prs.slides.add_slide(blank_layout)
            _add_header(s, prs, HEADER_TEXT_RIGHT_SEARCH_PC)
            _add_image_centered(s, slice_path, top_emu=650000)

    for slice_path in pc["slides"]:
        s = prs.slides.add_slide(blank_layout)
        _add_header(s, prs, HEADER_TEXT_RIGHT_PC)
        _add_image_centered(s, slice_path, top_emu=650000)

    if landing_pc and landing_pc.get("slides"):
        for slice_path in landing_pc["slides"]:
            s = prs.slides.add_slide(blank_layout)
            _add_header(s, prs, HEADER_TEXT_RIGHT_LANDING_PC)
            _add_image_centered(s, slice_path, top_emu=650000)

    # ---- MO: (검색결과) → 상세 5장 그리드 → 안내 5장 그리드 ----
    if search_mo and search_mo.get("slides"):
        for slice_path in search_mo["slides"]:
            s = prs.slides.add_slide(blank_layout)
            _add_header(s, prs, HEADER_TEXT_RIGHT_SEARCH_MO)
            _add_image_centered(s, slice_path, top_emu=650000)

    mo_slices = mo["slides"]
    for i in range(0, len(mo_slices), slides_per_mo_grid):
        chunk = mo_slices[i: i + slides_per_mo_grid]
        s = prs.slides.add_slide(blank_layout)
        _add_header(s, prs, HEADER_TEXT_RIGHT_MO)
        _add_image_grid(s, chunk, top_emu=700000)

    if landing_mo and landing_mo.get("slides"):
        landing_slices = landing_mo["slides"]
        for i in range(0, len(landing_slices), slides_per_mo_grid):
            chunk = landing_slices[i: i + slides_per_mo_grid]
            s = prs.slides.add_slide(blank_layout)
            _add_header(s, prs, HEADER_TEXT_RIGHT_LANDING_MO)
            _add_image_grid(s, chunk, top_emu=700000)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


# ---- CLI ------------------------------------------------------------------

DEFAULT_PC = (
    "https://card-search.naver.com/item?cardAdId=3257"
    "&query=%EC%8B%A0%ED%95%9C%EC%B9%B4%EB%93%9C%20Deep%20Once"
    "&cx=LNFvYD7ZLUB2W2S%2FZrj0TA%3D%3D"
)
DEFAULT_MO = (
    "https://m-card-search.naver.com/item?cardAdId=3257"
    "&query=%EC%8B%A0%ED%95%9C%EC%B9%B4%EB%93%9C%20Deep%20Once"
    "&cx=LNFvYD7ZLUB2W2S%2FZrj0TA%3D%3D"
)


def main():
    ap = argparse.ArgumentParser(description="신한카드 검색광고 게재보고 PPT 자동 생성")
    ap.add_argument("--pc", default=DEFAULT_PC, help="PC card-search.naver.com/item URL")
    ap.add_argument("--mo", default=DEFAULT_MO, help="MO m-card-search.naver.com/item URL")
    ap.add_argument("--search-pc", default=None, dest="search_pc",
                    help="(선택) PC 검색결과 URL (search.naver.com) — 카드 OneBox 캡처")
    ap.add_argument("--search-mo", default=None, dest="search_mo",
                    help="(선택) MO 검색결과 URL (m.search.naver.com) — 카드 OneBox 캡처")
    ap.add_argument("--landing", default=None,
                    help="(선택) 신한카드 자체 안내 페이지 URL — PC/MO 둘 다 캡처")
    ap.add_argument("--card", default="신한카드 Deep Once", help="카드명 (출력 파일명/타이틀용)")
    ap.add_argument(
        "--out",
        default=None,
        help="출력 PPT 경로 (생략 시 바탕화면에 자동 생성)",
    )
    ap.add_argument("--work", default=None, help="중간 캡처 저장 폴더 (생략 시 자동)")
    args = ap.parse_args()

    work = Path(args.work) if args.work else Path(tempfile.mkdtemp(prefix="report_"))
    out_default = Path(
        rf"C:\Users\FIN\Desktop\(핀플로우) 신한카드 멤버십 영업팀 신용카드 검색광고 게재보고_{args.card}.pptx"
    )
    out = Path(args.out) if args.out else out_default

    # Playwright sync API 는 한 process 에 한 인스턴스만 가능 → 순차 실행
    t0 = _time_mod.time()
    print(f"[1/2] 순차 캡처 시작 (work={work})", flush=True)
    pc = capture_pc(args.pc, work / "pc")
    print(f"      PC: {len(pc['slides'])}장 (H={pc['height']}px)", flush=True)
    mo = capture_mo(args.mo, work / "mo")
    print(f"      MO: {len(mo['slides'])}장 (H={mo['height']}px)", flush=True)

    search_pc_data = None
    search_mo_data = None
    landing_pc_data = None
    landing_mo_data = None
    if args.search_pc:
        search_pc_data = capture_search_onebox(args.search_pc, work / "search_pc", False)
        print(f"      Search PC: {len(search_pc_data['slides'])}장", flush=True)
    if args.search_mo:
        search_mo_data = capture_search_onebox(args.search_mo, work / "search_mo", True)
        print(f"      Search MO: {len(search_mo_data['slides'])}장", flush=True)
    if args.landing:
        landing_pc_data = capture_landing(args.landing, work / "landing_pc", False)
        print(f"      Landing PC: {len(landing_pc_data['slides'])}장 (H={landing_pc_data['height']})", flush=True)
        landing_mo_data = capture_landing(args.landing, work / "landing_mo", True)
        print(f"      Landing MO: {len(landing_mo_data['slides'])}장 (H={landing_mo_data['height']})", flush=True)
    print(f"      캡처 소요 {int(_time_mod.time() - t0)}s", flush=True)

    print("[2/2] PPT 생성 중...", flush=True)
    p = build_pptx(pc, mo, out, args.card,
                   search_pc=search_pc_data, search_mo=search_mo_data,
                   landing_pc=landing_pc_data, landing_mo=landing_mo_data)
    print(f"DONE → {p}")
    print(f"      (중간 캡처는 {work} 에 보존)")


if __name__ == "__main__":
    main()
